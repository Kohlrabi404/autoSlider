import collections
import collections.abc
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

import win32com.client as win32
import time
import subprocess

import cv2

class SlideShow :
    def __init__(self, path: str, amount: int):
        self.prs = Presentation()
        self.path = path
        self.amount = amount
    
    def make_slides(self, name: str) -> None :
        # Early return if file exists
        self.name = name
        if self.exist() :
            print("{}.pptx file already exists, maybe consider deleting current one.".format(self.name))
            return
        
        left = Inches(1)
        top = Inches(0)
        blank_slide_layout = self.prs.slide_layouts[6]
        for i in range(1, self.amount+1) :
            img_path = '{}\img_{:04}.png'.format(self.path, i)
            slide = self.prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(img_path, left, top, height = Inches(7.5))

        # Save the presentation
        self.prs.save("{}.pptx".format(name))

    def exist(self) -> bool :
        file_path = Path("{}.pptx".format(self.name))
        return file_path.is_file()

    def project(self, output_path, camera_index=0) -> None :

        self.output_path = output_path
        # Open the camera
        cap = cv2.VideoCapture(camera_index)

        # Check if the camera is opened successfully
        if not cap.isOpened():
            print("Failed to open the camera")
            return

        # Open the presentation in the default application
        subprocess.Popen(["start", "{}.pptx".format(self.name)], shell=True)
        time.sleep(5)

        # Connect to the running PowerPoint application
        powerpoint = win32.Dispatch("PowerPoint.Application")

        # Get the active presentation
        presentation = powerpoint.ActivePresentation

        # Start the slideshow
        slideshow_settings = presentation.SlideShowSettings
        slideshow_settings.Run()

        # Get the number of slides
        num_slides = presentation.Slides.Count
        print("Number of slides:", num_slides)

        for i in range(1, num_slides+1) :
            if i != 1 :
                presentation.SlideShowWindow.View.Next()
            #Capture and Save ith image
            ret, frame = cap.read()
            if not ret:
                print("Failed to capture frame")
                return
            cv2.imwrite("{}\img_{:04}.png".format(output_path, i), frame)
            time.sleep(2)

        # Close the PowerPoint application
        cap.release()
        powerpoint.Quit()

